# -*- coding: utf-8 -*-
"""
生成 系統架構報告_meeting.pptx
對應 系統架構報告_meeting.html 的 6 張投影片
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 顏色 ──
BLUE      = RGBColor(0x1A, 0x5F, 0xA8)
BLUE_L    = RGBColor(0x4A, 0x8F, 0xD4)
BLUE_BG   = RGBColor(0xE8, 0xF1, 0xFB)
GREEN     = RGBColor(0x4A, 0x7C, 0x4E)
GREEN_L   = RGBColor(0x7A, 0xB8, 0x7E)
GREEN_BG  = RGBColor(0xE6, 0xF4, 0xE7)
PURPLE    = RGBColor(0x6A, 0x4A, 0x9C)
PURPLE_L  = RGBColor(0x9A, 0x7A, 0xCC)
PURPLE_BG = RGBColor(0xF0, 0xEA, 0xF8)
ORANGE    = RGBColor(0xC0, 0x70, 0x20)
ORANGE_L  = RGBColor(0xE0, 0x98, 0x40)
ORANGE_BG = RGBColor(0xFD, 0xF3, 0xE3)
TEAL      = RGBColor(0x2A, 0x7A, 0x8A)
TEAL_BG   = RGBColor(0xE5, 0xF4, 0xF7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT      = RGBColor(0x1A, 0x1A, 0x1A)
MUTED     = RGBColor(0x5A, 0x62, 0x72)
BORDER    = RGBColor(0xD0, 0xD8, 0xE8)
BG        = RGBColor(0xF4, 0xF7, 0xFC)
RED_SOFT  = RGBColor(0xE5, 0x50, 0x50)
GREEN_OK  = RGBColor(0x20, 0x80, 0x48)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── 共用函式 ──
def S():
    return prs.slides.add_slide(BLANK)

def rect(slide, l, t, w, h, fill, line=None, lw=Pt(1), radius=False):
    shp_type = 5 if radius else 1
    s = slide.shapes.add_shape(shp_type, l, t, w, h)
    if radius:
        s.adjustments[0] = 0.04
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=Pt(14), color=TEXT, bold=False,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic
    return tb

def rtxt(shape, text, size=Pt(13), color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.color.rgb = color; r.font.bold = bold

def topbar(slide, c=BLUE):
    r = slide.shapes.add_shape(1, 0, 0, W, Pt(8))
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()

def botbar(slide):
    r = slide.shapes.add_shape(1, 0, H-Pt(6), W, Pt(6))
    r.fill.solid(); r.fill.fore_color.rgb = BLUE_L; r.line.fill.background()

def snum(slide, n, color=BLUE_L):
    txt(slide, f"{n:02d}", W-Inches(1.4), H-Inches(0.55), Inches(1.2), Inches(0.4),
        Pt(16), color, align=PP_ALIGN.RIGHT)

def tag(slide, text, l, t, fg, bg, border):
    s = slide.shapes.add_shape(5, l, t, Inches(4.5), Inches(0.42))
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = bg
    s.line.color.rgb = border; s.line.width = Pt(1.2)
    rtxt(s, text, Pt(13), fg, bold=True)

def slide_title(slide, text, em_text=""):
    full = text + em_text
    txt(slide, full, Inches(0.7), Inches(0.58), Inches(12.0), Inches(0.7),
        Pt(30), TEXT, bold=True)

def bullet_box(slide, l, t, w, h, title, title_color, items, bullet_color,
               bg=None, border=None, item_size=Pt(13)):
    if bg:
        rect(slide, l, t, w, h, bg, border, radius=True)
    txt(slide, title, l+Inches(0.2), t+Inches(0.12), w-Inches(0.4), Inches(0.38),
        Pt(15), title_color, bold=True)
    iy = t + Inches(0.55)
    for item in items:
        txt(slide, "▶  " + item, l+Inches(0.2), iy, w-Inches(0.35), Inches(0.38),
            item_size, MUTED, wrap=True)
        iy += Inches(0.42)

# ══════════════════════════════════════════════════════════
# SLIDE 1 — 系統架構總覽
# ══════════════════════════════════════════════════════════
slide = S()
rect(slide, 0, 0, W, H, BG)
topbar(slide); botbar(slide)

tag(slide, "Weekly Meeting  2026-04", Inches(0.7), Inches(0.12), BLUE, BLUE_BG, BLUE_L)
txt(slide, "半導體微影設備虛擬訓練系統 — 系統架構技術報告",
    Inches(5.5), Inches(0.1), Inches(7.5), Inches(0.55), Pt(22), TEXT, bold=True)

layers = [
    ("1", "Dataset  資料來源層",
     ["UCI SECOM 製程資料集", "1,567 wafer samples", "590 sensor features", "Pass 1,463 / Fail 104"],
     BLUE, BLUE_BG, BLUE_L),
    ("2", "User Interface  前端展示層",
     ["Three.js WebGL · Viewer.html", "3D 互動環境", "虛擬控制面板 HMI", "AI 對話 & 操作評估面板"],
     GREEN, GREEN_BG, GREEN_L),
    ("3", "AI Mentor System  AI 導師層",
     ["Qwen LLM 本地端推理", "AI Mentor 蘇格拉底引導", "Diagnostic / Operation Agent", "自適應四模式評分"],
     PURPLE, PURPLE_BG, PURPLE_L),
    ("4", "訓練流程層  Training Flow",
     ["理論學習 + 實機操作  整合式", "DUV 微影原理問答", "SOP 自主操作評估"],
     ORANGE, ORANGE_BG, ORANGE_L),
    ("5", "評分與問卷系統  Evaluation",
     ["操作正確性 + 問卷前後測", "三類受試者", "量化進步幅度統計"],
     TEAL, TEAL_BG, TEAL),
]
ly = Inches(0.82)
for num, name, chips, fc, bg_c, lc in layers:
    r = rect(slide, Inches(0.7), ly, Inches(12.0), Inches(1.02), bg_c, lc, radius=True)
    txt(slide, num, Inches(0.95), ly+Inches(0.28), Inches(0.4), Inches(0.45),
        Pt(24), fc, bold=True)
    txt(slide, name, Inches(1.5), ly+Inches(0.28), Inches(3.2), Inches(0.45),
        Pt(16), TEXT, bold=True)
    cx = Inches(5.0)
    for chip in chips:
        cw = Inches(1.9) if len(chip) <= 14 else Inches(2.8)
        cs = slide.shapes.add_shape(5, cx, ly+Inches(0.28), cw, Inches(0.42))
        cs.adjustments[0] = 0.3
        cs.fill.solid(); cs.fill.fore_color.rgb = WHITE
        cs.line.color.rgb = lc; cs.line.width = Pt(1)
        rtxt(cs, chip, Pt(12), fc, bold=True)
        cx += cw + Inches(0.18)
    ly += Inches(1.18)

snum(slide, 1)

# ══════════════════════════════════════════════════════════
# SLIDE 2 — Dataset
# ══════════════════════════════════════════════════════════
slide = S()
rect(slide, 0, 0, W, H, BG)
topbar(slide); botbar(slide)

tag(slide, "Layer 1  ·  Dataset", Inches(0.7), Inches(0.1), BLUE, BLUE_BG, BLUE_L)
txt(slide, "UCI SECOM 製程資料集  —  真實感測器統計基礎",
    Inches(5.4), Inches(0.08), Inches(7.7), Inches(0.55), Pt(24), TEXT, bold=True)

# Stat boxes
for i, (num, lbl, bg_c) in enumerate([
    ("1,567", "Wafer Samples", BLUE),
    ("590",   "Sensor Features", BLUE),
    ("1,463 / 104", "Pass / Fail", RED_SOFT),
]):
    bx = Inches(0.7) + i * Inches(2.3)
    s = rect(slide, bx, Inches(0.75), Inches(2.1), Inches(1.1), bg_c, radius=True)
    txt(slide, num, bx, Inches(0.85), Inches(2.1), Inches(0.55),
        Pt(28 if len(num)<=5 else 22), WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, lbl, bx, Inches(1.42), Inches(2.1), Inches(0.35),
        Pt(13), WHITE, align=PP_ALIGN.CENTER)

# Flow box
rect(slide, Inches(0.7), Inches(2.0), Inches(6.3), Inches(5.1), WHITE, BORDER, radius=True)
txt(slide, "前處理流程（SecomNoiseModel）", Inches(0.95), Inches(2.12),
    Inches(5.5), Inches(0.38), Pt(15), BLUE, bold=True)
flow_steps = [
    ("1", "填補 NaN（各欄中位數）", ""),
    ("2", "PCA 降維：590 特徵 → 前 20 主成分", "萃取與製程結果相關性最高的主成分軸"),
    ("3", "分離 Pass / Fail，建立雜訊統計分布", "正偏態雜訊模型（skewness ≈ 0.12）"),
    ("4", "計算 health_score（0 ~ 1）", "依 PCA 投影座標；分數越低代表狀態越差"),
]
fy = Inches(2.6)
for n, t1, t2 in flow_steps:
    s = rect(slide, Inches(0.95), fy, Inches(0.45), Inches(0.45), BLUE, radius=True)
    rtxt(s, n, Pt(12), WHITE, bold=True)
    txt(slide, t1, Inches(1.55), fy+Inches(0.02), Inches(5.2), Inches(0.35), Pt(14), TEXT)
    if t2:
        txt(slide, t2, Inches(1.55), fy+Inches(0.38), Inches(5.2), Inches(0.3), Pt(12), MUTED)
        fy += Inches(0.88)
    else:
        fy += Inches(0.6)

# Right usage cards
right_cards = [
    ("SECOM 在本系統中的角色", BLUE,
     "不預測良率，而是建立製程雜訊統計模型，驅動虛擬 HMI 的 590 個感測器即時讀值，使讀值符合真實統計分布"),
    ("CD 預測雜訊調變", BLUE,
     "正常狀態  σ = 1.5 nm     故障狀態  σ = 3.0 – 4.8 nm\nhealth_score 線性插值調變"),
    ("虛擬 HMI 面板呈現", BLUE,
     "590 個感測器即時讀值 / CDU Map / CD 趨勢圖（3σ 控制線）/ Overlay 向量場"),
    ("故障注入效果", BLUE,
     "注入 5 種故障後，noise σ 即時升高，反映至 CDU Map 及警報狀態"),
]
cy = Inches(0.75)
for title, tc, body in right_cards:
    rect(slide, Inches(7.3), cy, Inches(5.75), Inches(1.5), WHITE, BORDER, radius=True)
    txt(slide, title, Inches(7.55), cy+Inches(0.12), Inches(5.2), Inches(0.35), Pt(14), tc, bold=True)
    txt(slide, body, Inches(7.55), cy+Inches(0.52), Inches(5.2), Inches(0.85), Pt(13), MUTED, wrap=True)
    cy += Inches(1.68)

snum(slide, 2)

# ══════════════════════════════════════════════════════════
# SLIDE 3 — User Interface
# ══════════════════════════════════════════════════════════
slide = S()
rect(slide, 0, 0, W, H, BG)
topbar(slide); botbar(slide)

tag(slide, "Layer 2  ·  User Interface", Inches(0.7), Inches(0.1), GREEN, GREEN_BG, GREEN_L)
txt(slide, "Three.js WebGL  —  三大介面模組",
    Inches(5.7), Inches(0.08), Inches(7.5), Inches(0.55), Pt(24), TEXT, bold=True)

ui_cols = [
    ("3D 互動環境", "ASML DUV 虛擬機台", [
        "ASML DUV 高擬真 3D 模型（asml_duv.glb）",
        "PointerLock 第一人稱視角控制",
        "WASD 移動 / 滑鼠視角旋轉",
        "Raycaster 偵測鄰近零件，按 [E] 觸發互動",
        "故障發生時對應 Mesh 橘色發光警示",
        "11 個可互動零件群組（鏡組、載台、冷卻…）",
    ], None),
    ("虛擬控制面板", "HMI 感測器即時顯示", [
        "590 個感測器讀值（SECOM 驅動）",
        "Dose / Focus / LensTemp / CDU 3σ / Overlay",
        "CDU Map 2D 晶圓空間分布圖",
        "CD 趨勢圖（含 ±3σ 控制線）",
        "Overlay 向量場可視化",
        "故障觸發時自動顯示紅色警報欄",
    ], "GET /api/hmi  →  JSON 感測器快照"),
    ("AI 對話 & 操作評估面板", "按 [C] 開啟 / 按 [E] 互動", [
        "[C] 鍵開啟 AI Chat 對話框，即時蘇格拉底問答",
        "靠近零件按 [E]：出現「檢查」「操作」按鈕",
        "學員自行判斷操作，系統即時評估對錯",
        "右上角即時分數顯示（滿分 100）",
        "求助學長（扣 10 分）→ 差異化提示",
    ], "POST /api/action  →  SOP 比對  →  即時回饋"),
]
cx = Inches(0.55)
cw = Inches(4.05)
for title, sub, items, api in ui_cols:
    rect(slide, cx, Inches(0.72), cw, Inches(6.45), GREEN_BG, GREEN_L, radius=True)
    txt(slide, title, cx+Inches(0.22), Inches(0.85), cw-Inches(0.3), Inches(0.45),
        Pt(18), TEXT, bold=True)
    txt(slide, sub, cx+Inches(0.22), Inches(1.32), cw-Inches(0.3), Inches(0.32),
        Pt(13), GREEN)
    iy = Inches(1.75)
    for item in items:
        txt(slide, "▶  " + item, cx+Inches(0.22), iy, cw-Inches(0.35), Inches(0.38),
            Pt(13), MUTED, wrap=True)
        iy += Inches(0.43)
    if api:
        s = rect(slide, cx+Inches(0.18), Inches(6.65), cw-Inches(0.35), Inches(0.42),
                 GREEN, radius=True)
        rtxt(s, api, Pt(12), WHITE, bold=True)
    cx += cw + Inches(0.28)

snum(slide, 3)

# ══════════════════════════════════════════════════════════
# SLIDE 4 — AI Mentor System
# ══════════════════════════════════════════════════════════
slide = S()
rect(slide, 0, 0, W, H, BG)
topbar(slide); botbar(slide)

tag(slide, "Layer 3  ·  AI Mentor System", Inches(0.7), Inches(0.1), PURPLE, PURPLE_BG, PURPLE_L)
txt(slide, "Qwen LLM 本地端推理  —  三子系統協同",
    Inches(5.7), Inches(0.08), Inches(7.5), Inches(0.55), Pt(24), TEXT, bold=True)

col_w = Inches(4.1)
col_tops = [Inches(0.72), Inches(0.72), Inches(0.72)]
col_ls   = [Inches(0.55), Inches(4.75), Inches(8.95)]
col_bgs  = [PURPLE_BG, RGBColor(0xF0,0xF4,0xFF), BLUE_BG]
col_lcs  = [PURPLE_L,  RGBColor(0x88,0x88,0xDD), BLUE_L]
col_tc   = [PURPLE,    RGBColor(0x40,0x40,0xBB), BLUE]

# Card 1: AI Mentor
rect(slide, col_ls[0], col_tops[0], col_w, Inches(6.45), col_bgs[0], col_lcs[0], radius=True)
txt(slide, "AI Mentor", col_ls[0]+Inches(0.22), Inches(0.82), Inches(3.6), Inches(0.42),
    Pt(20), col_tc[0], bold=True)
txt(slide, "AIScenarioMentor · ProactiveMentor",
    col_ls[0]+Inches(0.22), Inches(1.26), Inches(3.6), Inches(0.3), Pt(12), MUTED)
items1 = [
    ("蘇格拉底式對話引導",
     "不直接給答案，用追問引導主動推理；依場景脈絡生成差異化問題"),
    ("主動故障告警生成",
     "故障注入後主動觸發問題，5 種故障 × 3 個隨機變體，避免背答案"),
    ("雙層語義評估",
     "同義詞關鍵字群組（偏移/shift/歪掉）+ LLM 語義評估（+5/+2/−3）"),
]
iy = Inches(1.65)
for t, d in items1:
    rect(slide, col_ls[0]+Inches(0.18), iy, col_w-Inches(0.35), Inches(1.35),
         WHITE, BORDER, radius=True)
    txt(slide, t, col_ls[0]+Inches(0.38), iy+Inches(0.1), Inches(3.4), Inches(0.35),
        Pt(14), TEXT, bold=True)
    txt(slide, d, col_ls[0]+Inches(0.38), iy+Inches(0.48), Inches(3.4), Inches(0.75),
        Pt(12), MUTED, wrap=True)
    iy += Inches(1.5)
s = rect(slide, col_ls[0]+Inches(0.18), Inches(6.48), col_w-Inches(0.35), Inches(0.55),
         PURPLE, radius=True)
rtxt(s, "本地端 Qwen LLM  完全不連網  符合晶圓廠資安", Pt(13), WHITE, bold=True)

# Card 2: Agent
rect(slide, col_ls[1], col_tops[1], col_w, Inches(6.45), col_bgs[1], col_lcs[1], radius=True)
txt(slide, "Diagnostic / Operation Agent",
    col_ls[1]+Inches(0.22), Inches(0.82), Inches(3.6), Inches(0.42), Pt(18), col_tc[1], bold=True)
txt(slide, "A2ACoordinator 協調",
    col_ls[1]+Inches(0.22), Inches(1.26), Inches(3.6), Inches(0.3), Pt(12), MUTED)
items2 = [
    ("故障根因分析 (Diagnostic Agent)",
     "根據感測器異常值推斷故障類型，結合 SECOM health_score 輔助診斷"),
    ("SOP 操作管理 (Operation Agent)",
     "維護 sop_definitions.py 中 5 種故障的有序步驟序列，逐步比對學員操作"),
    ("錯誤三級分類",
     "partial_action / partial_component / full_wrong，對應不同扣分"),
    ("Safety Agent",
     "PPE 穿戴合規提醒，佔評分 30%"),
]
iy = Inches(1.65)
for t, d in items2:
    rect(slide, col_ls[1]+Inches(0.18), iy, col_w-Inches(0.35), Inches(1.1),
         WHITE, BORDER, radius=True)
    txt(slide, t, col_ls[1]+Inches(0.38), iy+Inches(0.08), Inches(3.4), Inches(0.35),
        Pt(13), TEXT, bold=True)
    txt(slide, d, col_ls[1]+Inches(0.38), iy+Inches(0.45), Inches(3.4), Inches(0.58),
        Pt(12), MUTED, wrap=True)
    iy += Inches(1.2)

# Card 3: Adaptive
rect(slide, col_ls[2], col_tops[2], col_w, Inches(6.45), col_bgs[2], col_lcs[2], radius=True)
txt(slide, "自適應四模式評分",
    col_ls[2]+Inches(0.22), Inches(0.82), Inches(3.6), Inches(0.42), Pt(20), col_tc[2], bold=True)
txt(slide, "ActionSession  起始 100 分",
    col_ls[2]+Inches(0.22), Inches(1.26), Inches(3.6), Inches(0.3), Pt(12), MUTED)
modes = [
    ("> 85 分",  "挑戰  Challenge",  RGBColor(0x1A,0x2A,0x3E), WHITE,   "不給提示，引導自行分析"),
    ("64–85 分", "標準  Standard",   BLUE_BG,                  BLUE,    "指出讀值偏差與應前往部件"),
    ("40–63 分", "鷹架  Scaffolding",ORANGE_BG,                ORANGE,  "給部件名稱 + 操作方向"),
    ("< 40 分",  "補救  Remedial",   RGBColor(0xFD,0xE8,0xE8), RGBColor(0xB0,0x30,0x30), "完整步驟：部件、動作、原因"),
]
iy = Inches(1.65)
for score, name, bg_c, fc, desc in modes:
    rect(slide, col_ls[2]+Inches(0.18), iy, col_w-Inches(0.35), Inches(1.05),
         bg_c, None, radius=True)
    txt(slide, name, col_ls[2]+Inches(0.38), iy+Inches(0.08), Inches(2.0), Inches(0.38),
        Pt(14), fc, bold=True)
    txt(slide, score, col_ls[2]+Inches(2.5), iy+Inches(0.08), Inches(1.3), Inches(0.38),
        Pt(16), fc, bold=True)
    txt(slide, desc, col_ls[2]+Inches(0.38), iy+Inches(0.5), Inches(3.4), Inches(0.45),
        Pt(12), MUTED)
    iy += Inches(1.2)
rect(slide, col_ls[2]+Inches(0.18), Inches(6.48), col_w-Inches(0.35), Inches(0.42),
     WHITE, BORDER, radius=True)
txt(slide, "求助學長：扣 10 分後依當前分數給差異化方向提示",
    col_ls[2]+Inches(0.35), Inches(6.54), Inches(3.5), Inches(0.3), Pt(12), MUTED)

snum(slide, 4)

# ══════════════════════════════════════════════════════════
# SLIDE 5 — Training Flow
# ══════════════════════════════════════════════════════════
slide = S()
rect(slide, 0, 0, W, H, BG)
topbar(slide); botbar(slide)

tag(slide, "Layer 4  ·  Training Flow", Inches(0.7), Inches(0.1), ORANGE, ORANGE_BG, ORANGE_L)
txt(slide, "整合式訓練設計  —  理論與實機同步進行",
    Inches(5.7), Inches(0.08), Inches(7.5), Inches(0.55), Pt(24), TEXT, bold=True)

half_w = Inches(5.9)

# Theory col
rect(slide, Inches(0.55), Inches(0.72), half_w, Inches(6.45), ORANGE_BG, ORANGE_L, radius=True)
s_hdr = rect(slide, Inches(0.55), Inches(0.72), half_w, Inches(0.6), ORANGE, radius=True)
rtxt(s_hdr, "理論學習", Pt(20), WHITE, bold=True)
theory_items = [
    "DUV 微影原理（Rayleigh、Bossung Curve、DOF）",
    "製程參數關聯：Dose / Focus / NA / CDU",
    "物理模型：鏡片熱方程式（雙指數）、Overlay 誤差來源",
    "AI 蘇格拉底式問答鞏固——不直接說答案，用追問引導",
    "知識點追蹤（KnowledgeTracker）：thermal / vacuum / optical 六類",
]
iy = Inches(1.45)
for item in theory_items:
    txt(slide, "▶  " + item, Inches(0.8), iy, Inches(5.4), Inches(0.42),
        Pt(14), MUTED, wrap=True)
    iy += Inches(0.46)

rect(slide, Inches(0.72), Inches(5.3), half_w-Inches(0.35), Inches(1.58),
     WHITE, BORDER, radius=True)
txt(slide, "整合式設計：無分離測驗門檻",
    Inches(0.95), Inches(5.42), Inches(5.2), Inches(0.38), Pt(15), TEXT, bold=True)
txt(slide, "AI 導師根據學員輸入內容，動態切換「理論引導」或「實作指引」，兩種模式在同一對話框無縫交替",
    Inches(0.95), Inches(5.85), Inches(5.2), Inches(0.9), Pt(13), MUTED, wrap=True)

# Arrow
txt(slide, "⇄", Inches(6.55), Inches(3.3), Inches(0.7), Inches(0.7),
    Pt(40), MUTED, align=PP_ALIGN.CENTER)

# Practice col
rect(slide, Inches(7.35), Inches(0.72), half_w, Inches(6.45), BLUE_BG, BLUE_L, radius=True)
s_hdrp = rect(slide, Inches(7.35), Inches(0.72), half_w, Inches(0.6), BLUE, radius=True)
rtxt(s_hdrp, "實機操作（自主評估）", Pt(20), WHITE, bold=True)
prac_items = [
    "5 種故障情境：鏡片熱點 / 光罩污染 / 載台誤差 / Dose 漂移 / Focus 漂移",
    "學員自行判斷靠近哪個零件、選擇何種操作（非逐步指引）",
    "POST /api/action → sop_definitions.py 比對 → 即時回饋",
    "連續答對 ≥ 3 次 → 升難度；連續答錯 ≥ 3 次 → 降難度",
    "感測器 HMI 即時反映故障狀態（SECOM 驅動 + 物理模型）",
]
iy = Inches(1.45)
for item in prac_items:
    txt(slide, "▶  " + item, Inches(7.6), iy, Inches(5.4), Inches(0.42),
        Pt(14), MUTED, wrap=True)
    iy += Inches(0.46)

rect(slide, Inches(7.52), Inches(5.3), half_w-Inches(0.35), Inches(1.58),
     WHITE, BLUE_L, radius=True)
txt(slide, "5 種 SOP，各 4–5 個有序步驟",
    Inches(7.75), Inches(5.42), Inches(5.2), Inches(0.38), Pt(15), BLUE, bold=True)
txt(slide, "錯誤依三級分類扣分（2–15 分），求助扣 10 分，全程 100 分制動態計分",
    Inches(7.75), Inches(5.85), Inches(5.2), Inches(0.9), Pt(13), MUTED, wrap=True)

snum(slide, 5)

# ══════════════════════════════════════════════════════════
# SLIDE 6 — Evaluation
# ══════════════════════════════════════════════════════════
slide = S()
rect(slide, 0, 0, W, H, BG)
topbar(slide); botbar(slide)

tag(slide, "Layer 5  ·  Evaluation", Inches(0.7), Inches(0.1), TEAL, TEAL_BG, TEAL)
txt(slide, "評分與問卷系統  —  量化能力驗證",
    Inches(5.4), Inches(0.08), Inches(7.8), Inches(0.55), Pt(24), TEXT, bold=True)

half_w2 = Inches(6.15)

# Left: Operation scoring
rect(slide, Inches(0.55), Inches(0.72), half_w2, Inches(6.45), TEAL_BG, TEAL, radius=True)
txt(slide, "操作正確性評分",
    Inches(0.8), Inches(0.85), Inches(5.5), Inches(0.45), Pt(22), TEAL, bold=True)
ops_items = [
    ("故障類型判斷（Diagnostic）",
     "依感測器讀值正確辨識故障類型（lens_hotspot / contamination…）"),
    ("SOP 操作順序（Operation）",
     "按照 sop_definitions.py 有序步驟執行；部分匹配依類型扣 2–15 分"),
    ("理論回答正確性",
     "同義詞群組關鍵字 +7/+3 分；LLM 語義分析 +5/+2/−3 分"),
    ("安全合規（Safety）",
     "PPE 穿戴 / 鎖定程序 / 操作規範遵守，Safety Agent 即時稽核"),
]
iy = Inches(1.45)
for title, desc in ops_items:
    rect(slide, Inches(0.75), iy, half_w2-Inches(0.4), Inches(1.15),
         WHITE, BORDER, radius=True)
    txt(slide, title, Inches(1.0), iy+Inches(0.1), Inches(5.5), Inches(0.38),
        Pt(14), TEXT, bold=True)
    txt(slide, desc, Inches(1.0), iy+Inches(0.52), Inches(5.5), Inches(0.55),
        Pt(13), MUTED, wrap=True)
    iy += Inches(1.3)

# Right: Quiz
rect(slide, Inches(7.0), Inches(0.72), half_w2, Inches(6.45), BLUE_BG, BLUE_L, radius=True)
txt(slide, "問卷前後測",
    Inches(7.25), Inches(0.85), Inches(5.5), Inches(0.45), Pt(22), BLUE, bold=True)

# Group pills
txt(slide, "受試者分組：", Inches(7.25), Inches(1.42), Inches(5.5), Inches(0.35), Pt(14), TEXT, bold=True)
for i, (gtext, gc) in enumerate([
    ("無半導體背景",     RGBColor(0xE8,0xF1,0xFB)),
    ("工程師 0–5 年",   RGBColor(0xF0,0xEA,0xF8)),
    ("資深工程師 5 年+", RGBColor(0xFD,0xF3,0xE3)),
]):
    gx = Inches(7.25) + i * Inches(1.95)
    s = rect(slide, gx, Inches(1.82), Inches(1.85), Inches(0.42), gc, radius=True)

for i, (gtext, gc, fc) in enumerate([
    ("無半導體背景",     RGBColor(0xE8,0xF1,0xFB), BLUE),
    ("工程師 0–5 年",   RGBColor(0xF0,0xEA,0xF8), PURPLE),
    ("資深工程師 5 年+", RGBColor(0xFD,0xF3,0xE3), ORANGE),
]):
    gx = Inches(7.25) + i * Inches(1.95)
    txt(slide, gtext, gx+Inches(0.1), Inches(1.87), Inches(1.7), Inches(0.3),
        Pt(12), fc, bold=True, align=PP_ALIGN.CENTER)

# Test flow card
rect(slide, Inches(7.2), Inches(2.4), half_w2-Inches(0.25), Inches(1.1), WHITE, BORDER, radius=True)
txt(slide, "測試流程", Inches(7.45), Inches(2.52), Inches(5.4), Inches(0.35), Pt(14), TEXT, bold=True)
txt(slide, "前測  →  系統操作訓練（20–30 分鐘）→  後測\n概念理解問卷，滿分 15 分",
    Inches(7.45), Inches(2.92), Inches(5.4), Inches(0.5), Pt(13), MUTED)

# Best result card
rect(slide, Inches(7.2), Inches(3.65), half_w2-Inches(0.25), Inches(1.2), WHITE, BORDER, radius=True)
txt(slide, "最高成效族群：0–5 年新進工程師",
    Inches(7.45), Inches(3.77), Inches(5.4), Inches(0.38), Pt(14), TEXT, bold=True)
txt(slide, "前測 11.67  →  後測 13.33，進步 11.13%\n具基礎但尚未熟練者受益最大",
    Inches(7.45), Inches(4.18), Inches(5.4), Inches(0.6), Pt(13), MUTED)

# Result box
s_res = rect(slide, Inches(7.2), Inches(5.05), half_w2-Inches(0.25), Inches(1.82),
             BLUE, radius=True)
txt(slide, "+7.27%", Inches(7.4), Inches(5.18), Inches(2.0), Inches(0.9),
    Pt(44), WHITE, bold=True)
txt(slide, "整體平均進步幅度",   Inches(9.45), Inches(5.18), Inches(3.5), Inches(0.35), Pt(14), WHITE, bold=True)
txt(slide, "前測 10.45  →  後測 11.55 分", Inches(9.45), Inches(5.55), Inches(3.5), Inches(0.32), Pt(13), WHITE)
txt(slide, "平均進步 +1.09 分（滿分 15）",  Inches(9.45), Inches(5.88), Inches(3.5), Inches(0.32), Pt(13), WHITE)

snum(slide, 6)

# ── 儲存 ──
out = r"c:\Users\user\Desktop\在職碩\OneDrive - 長庚大學\長庚碩班\論文\系統架構報告_meeting.pptx"
prs.save(out)
print("Saved: " + out)
