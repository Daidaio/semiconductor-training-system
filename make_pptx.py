# -*- coding: utf-8 -*-
"""
生成口試報告 .pptx
依 口試報告.html 的內容重建為 PowerPoint 格式
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── 顏色定義（對應 CSS oklch 轉換後的 RGB）──
ACCENT      = RGBColor(0x0C, 0x5F, 0x8E)   # oklch(0.42 0.18 225) ≈ #0C5F8E
ACCENT_LIGHT= RGBColor(0x5E, 0xA8, 0xD4)   # oklch(0.72 0.12 225) ≈ #5EA8D4
ACCENT_BG   = RGBColor(0xEB, 0xF4, 0xFA)   # oklch(0.96 0.04 225)
GOLD        = RGBColor(0xB5, 0x8A, 0x2F)   # oklch(0.65 0.14 75)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT        = RGBColor(0x1A, 0x1A, 0x1A)
MUTED       = RGBColor(0x6B, 0x72, 0x80)
BG_LIGHT    = RGBColor(0xF8, 0xFA, 0xFC)
BORDER      = RGBColor(0xE5, 0xE7, 0xEB)
RED_SOFT    = RGBColor(0xF5, 0xD0, 0xC8)
GREEN_SOFT  = RGBColor(0xC8, 0xED, 0xD9)

W = Inches(13.33)  # 1920px → 16:9 widescreen
H = Inches(7.5)    # 1080px

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # Blank layout

def add_slide():
    return prs.slides.add_slide(blank_layout)

def rgb(r, g, b):
    return RGBColor(r, g, b)

def top_bar(slide, color=ACCENT):
    bar = slide.shapes.add_shape(1, 0, 0, W, Pt(6))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def bottom_bar(slide):
    bar = slide.shapes.add_shape(1, 0, H - Pt(6), W, Pt(6))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_LIGHT
    bar.line.fill.background()

def slide_num(slide, n, color=ACCENT_LIGHT):
    txb = slide.shapes.add_textbox(W - Inches(1.5), H - Inches(0.7), Inches(1.2), Inches(0.4))
    tf = txb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    run = p.add_run(); run.text = f"{n:02d}"
    run.font.size = Pt(16); run.font.color.rgb = color

def add_label(slide, text, left, top, width=Inches(4), height=Inches(0.4)):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text.upper()
    run.font.size = Pt(13); run.font.color.rgb = ACCENT
    run.font.bold = True

def add_title(slide, text, left, top, width, height, size=Pt(36), color=TEXT, bold=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text
    run.font.size = size; run.font.color.rgb = color; run.font.bold = bold
    return txb

def add_text(slide, text, left, top, width, height, size=Pt(14), color=TEXT, bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = size; run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic
    return txb

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, text=None, text_size=Pt(13), text_color=TEXT):
    shape = slide.shapes.add_shape(5, left, top, width, height)  # 5 = rounded rectangle
    shape.adjustments[0] = 0.05
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = text
        run.font.size = text_size; run.font.color.rgb = text_color; run.font.bold = True
    return shape

# ═══════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════
slide = add_slide()

# Left half background (white)
add_rect(slide, 0, 0, W/2, H, WHITE)
# Right half background (accent blue)
add_rect(slide, W/2, 0, W/2, H, ACCENT)

top_bar(slide)
bottom_bar(slide)

# Tag pill
add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(3.2), Inches(0.45),
                 ACCENT_BG, ACCENT, "長庚大學 碩士論文口試", Pt(13), ACCENT)

# Main title
add_title(slide, "基於數位孿生之", Inches(0.8), Inches(2.1), Inches(5.8), Inches(0.7), Pt(34), TEXT)
add_title(slide, "曝光機故障診斷訓練系統", Inches(0.8), Inches(2.8), Inches(6.0), Inches(0.8), Pt(34), ACCENT)

# Subtitle
add_text(slide, "Digital Twin-based Training System\nfor Lithography Fault Diagnosis",
         Inches(0.8), Inches(3.8), Inches(5.8), Inches(0.9),
         Pt(15), MUTED)

# Meta rows
meta_top = Inches(5.0)
for label, val in [("研究生", "趙威豪 M1321114"), ("指導教授", "張永華 教授"),
                   ("單位", "長庚大學 電機系"), ("日期", "2026 年 4 月")]:
    add_text(slide, label, Inches(0.8), meta_top, Inches(1.0), Inches(0.35), Pt(13), MUTED)
    add_text(slide, val,   Inches(1.9), meta_top, Inches(3.5), Inches(0.35), Pt(13), TEXT, bold=True)
    meta_top += Inches(0.4)

# Right side keyword pills
pills = ["數位孿生  Digital Twin",
         "ProactiveMentor AI 導師",
         "自適應 SOP 評估引擎",
         "生成式 AI / 蘇格拉底追問"]
pill_top = Inches(3.2)
for pill in pills:
    p_shape = slide.shapes.add_shape(5, Inches(7.2), pill_top, Inches(5.5), Inches(0.5))
    p_shape.adjustments[0] = 0.5
    p_shape.fill.solid(); p_shape.fill.fore_color.rgb = RGBColor(0x20, 0x60, 0x80)
    p_shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); p_shape.line.width = Pt(1)
    tf = p_shape.text_frame; tf.word_wrap = False
    par = tf.paragraphs[0]; par.alignment = PP_ALIGN.CENTER
    r = par.add_run(); r.text = pill
    r.font.size = Pt(14); r.font.color.rgb = WHITE
    pill_top += Inches(0.65)

# Machine label
add_text(slide, "ASML TWINSCAN NXT:870", Inches(7.5), Inches(5.8), Inches(5.0), Inches(0.5),
         Pt(16), WHITE, bold=True, align=PP_ALIGN.CENTER)

slide_num(slide, 1, RGBColor(0xFF, 0xFF, 0xFF))

# ═══════════════════════════════════════════════════════════
# SLIDE 2: BACKGROUND
# ═══════════════════════════════════════════════════════════
slide = add_slide()
top_bar(slide); bottom_bar(slide)

add_label(slide, "研究背景", Inches(0.8), Inches(0.55))
add_title(slide, "微影製程：半導體製造的心臟",
          Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.7), Pt(32))

# Left stats
stats = [("30–35%", "微影製程佔整體晶圓廠製造成本比例"),
         ("數千萬", "ASML DUV 機台單台採購成本（美元）"),
         ("$10萬+", "機台停機一小時的機會成本（美元）")]
stat_top = Inches(1.9)
for num, desc in stats:
    add_text(slide, num, Inches(0.8), stat_top, Inches(5.5), Inches(0.85),
             Pt(44), ACCENT, bold=True)
    add_text(slide, desc, Inches(0.8), stat_top + Inches(0.85), Inches(5.5), Inches(0.45),
             Pt(14), MUTED)
    stat_top += Inches(1.6)

# Right fact cards
facts = [("技術門檻極高", "操作人員需具備光學物理、製程控制、故障診斷的深厚跨域知識，培訓週期漫長"),
         ("新進工程師缺乏安全練習場域", "任何操作失誤都可能觸發機台聯鎖停機、批量晶圓報廢，代價極為昂貴"),
         ("稀有故障難以演練", "光罩污染、晶圓台誤差等高危故障極少發生，但處置不當後果嚴重，新人缺乏實際經驗")]
fact_top = Inches(1.9)
for title, body in facts:
    r = add_rect(slide, Inches(6.8), fact_top, Inches(6.2), Inches(1.6), ACCENT_BG, ACCENT_LIGHT)
    bar = add_rect(slide, Inches(6.8), fact_top, Pt(5), Inches(1.6), ACCENT)
    add_text(slide, title, Inches(7.1), fact_top + Inches(0.15), Inches(5.8), Inches(0.4),
             Pt(15), TEXT, bold=True)
    add_text(slide, body, Inches(7.1), fact_top + Inches(0.55), Inches(5.8), Inches(0.9),
             Pt(13), MUTED, wrap=True)
    fact_top += Inches(1.8)

slide_num(slide, 2)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: PAIN POINTS
# ═══════════════════════════════════════════════════════════
slide = add_slide()
top_bar(slide); bottom_bar(slide)

add_label(slide, "問題定義", Inches(0.8), Inches(0.55))
add_title(slide, "傳統訓練的四大痛點",
          Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.7), Pt(32))

pains = [
    ("依賴實機機時", "新進工程師需 2–4 週實機訓練，佔用昂貴設備時間，排程困難", "⚠ 設備使用率降低，影響產能排程"),
    ("資深工程師全程陪同", "必須消耗寶貴的資深人力資源，無法同時處理其他技術任務", "⚠ 人力成本高，資深人才難以善用"),
    ("失誤後果嚴重", "劑量設定錯誤、光罩裝卸失誤→聯鎖停機→批量晶圓報廢", "⚠ 真實代價難以承受，壓力大"),
    ("稀有故障無法練習", "高危故障極少發生，新人缺乏診斷與處置的實際操作機會", "⚠ 關鍵技能只靠理論，上線後風險高"),
]
card_left = Inches(0.8)
for title, body, impact in pains:
    add_rounded_rect(slide, card_left, Inches(1.9), Inches(2.9), Inches(5.2), BG_LIGHT, BORDER)
    add_text(slide, title, card_left + Inches(0.2), Inches(2.3), Inches(2.6), Inches(0.5),
             Pt(16), TEXT, bold=True)
    add_text(slide, body, card_left + Inches(0.2), Inches(2.9), Inches(2.6), Inches(1.5),
             Pt(13), MUTED, wrap=True)
    add_rounded_rect(slide, card_left + Inches(0.15), Inches(6.0), Inches(2.6), Inches(0.7),
                     RGBColor(0xFE, 0xF0, 0xEE), RGBColor(0xF5, 0xB5, 0xA0),
                     impact, Pt(12), RGBColor(0x8B, 0x35, 0x22))
    card_left += Inches(3.15)

slide_num(slide, 3)

# ═══════════════════════════════════════════════════════════
# SLIDE 4: SOLUTION OVERVIEW (blue background)
# ═══════════════════════════════════════════════════════════
slide = add_slide()
add_rect(slide, 0, 0, W, H, ACCENT)

add_title(slide, "數位孿生訓練平台\n讓工程師在零風險環境中習得故障診斷",
          Inches(1.5), Inches(0.8), Inches(10.5), Inches(1.6), Pt(32), WHITE)

trio = [
    ("01", "3D 虛擬訓練環境\nDigital Twin Environment", "Vue.js + Three.js WebGL 第一人稱 3D\n依 ASML TWINSCAN NXT:870 建模"),
    ("02", "自適應 SOP 評估引擎\nActionSession", "100 分制動態評分，四級提示層級\n（Challenge/Standard/Scaffolding/Remedial）"),
    ("03", "ProactiveMentor AI 導師\nLocal Qwen LLM", "蘇格拉底式追問，雙層語義評估\n（同義詞關鍵字 + LLM）"),
]
card_l = Inches(0.6)
for num, label, sub in trio:
    shape = slide.shapes.add_shape(5, card_l, Inches(2.8), Inches(4.0), Inches(3.2))
    shape.adjustments[0] = 0.05
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x18, 0x5F, 0x8A)
    shape.line.color.rgb = RGBColor(0x50, 0x90, 0xBB); shape.line.width = Pt(1.5)
    add_text(slide, num, card_l + Inches(0.25), Inches(3.0), Inches(0.8), Inches(0.7),
             Pt(36), RGBColor(0xA0, 0xC8, 0xE0), bold=True)
    add_text(slide, label, card_l + Inches(0.25), Inches(3.75), Inches(3.5), Inches(0.8),
             Pt(16), WHITE, bold=True)
    add_text(slide, sub, card_l + Inches(0.25), Inches(4.65), Inches(3.5), Inches(1.0),
             Pt(13), RGBColor(0xC0, 0xD8, 0xEA), wrap=True)
    card_l += Inches(4.3)

pills4 = ["ArF 193 nm 浸潤式 DUV", "UCI SECOM 資料集驅動", "5 種故障情境", "完全本地端推理·資安符合"]
pl = Inches(0.8)
for p4 in pills4:
    pw = Inches(2.8) if len(p4) < 12 else Inches(3.5)
    sh = slide.shapes.add_shape(5, pl, Inches(6.45), pw, Inches(0.6))
    sh.adjustments[0] = 0.5
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x20, 0x60, 0x8A)
    sh.line.color.rgb = RGBColor(0x70, 0xA8, 0xCC); sh.line.width = Pt(1)
    tf = sh.text_frame; par = tf.paragraphs[0]; par.alignment = PP_ALIGN.CENTER
    r = par.add_run(); r.text = p4; r.font.size = Pt(14); r.font.color.rgb = WHITE
    pl += pw + Inches(0.3)

slide_num(slide, 4, RGBColor(0xFF, 0xFF, 0xFF))

# ═══════════════════════════════════════════════════════════
# SLIDE 5: ARCHITECTURE
# ═══════════════════════════════════════════════════════════
slide = add_slide()
top_bar(slide); bottom_bar(slide)

add_label(slide, "系統設計 · 圖 1 系統整體架構圖", Inches(0.8), Inches(0.55), Inches(8))
add_title(slide, "三層式架構：渲染 / 邏輯 / AI 推理解耦",
          Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.7), Pt(32))

layers = [
    ("L1", "前端展示層  Frontend Layer", ACCENT_BG, ACCENT_LIGHT,
     "Vue.js  ·  Three.js (WebGL)  ·  第一人稱 3D 虛擬環境  ·  PointerLockControls  ·  操作行為記錄"),
    ("L2", "後端邏輯層  Backend Layer", BG_LIGHT, BORDER,
     "Python FastAPI  ·  RESTful API  ·  自適應 SOP 評估引擎  ·  即時得分計算  ·  訓練狀態資料庫"),
    ("L3", "AI 推理層  AI Inference Layer", ACCENT_BG, ACCENT_LIGHT,
     "本地端 Qwen LLM  ·  ProactiveMentor  ·  同義詞群組評估  ·  LLM 語義分析  ·  不連網外部 API"),
]
ly = Inches(2.0)
for num, name, bg, line_c, chips in layers:
    add_rounded_rect(slide, Inches(0.8), ly, Inches(11.8), Inches(1.5), bg, line_c)
    add_text(slide, num, Inches(1.1), ly + Inches(0.1), Inches(0.7), Inches(0.4),
             Pt(18), ACCENT_LIGHT, bold=True)
    add_text(slide, name, Inches(1.9), ly + Inches(0.1), Inches(5.5), Inches(0.4),
             Pt(16), TEXT, bold=True)
    add_text(slide, chips, Inches(1.1), ly + Inches(0.65), Inches(10.8), Inches(0.6),
             Pt(13), MUTED)
    ly += Inches(1.75)

slide_num(slide, 5)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: DIGITAL TWIN
# ═══════════════════════════════════════════════════════════
slide = add_slide()
top_bar(slide); bottom_bar(slide)

add_label(slide, "核心技術 1 · 圖 2 虛擬訓練環境實機截圖", Inches(0.8), Inches(0.55), Inches(9))
add_title(slide, "3D 虛擬訓練環境（數位孿生）",
          Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.6), Pt(32))

add_text(slide, "依 ASML TWINSCAN NXT:870 實機尺寸比例建模",
         Inches(0.8), Inches(1.65), Inches(7.5), Inches(0.4), Pt(14), ACCENT, bold=True)

subsys = [
    ("雷射光源 ArF DUV Laser", "193 nm 浸潤式深紫外光"),
    ("投影鏡組 Projection Optics", "4× 縮小投影，NA 控制"),
    ("照明系統 Illumination", "均化光束，σ 相干度設定"),
    ("晶圓載台 Wafer Stage", "nm 等級步進定位"),
]
sy = Inches(2.2)
sx = Inches(0.8)
for i, (name, desc) in enumerate(subsys):
    cx = sx if i % 2 == 0 else sx + Inches(3.7)
    cy = sy if i < 2 else sy + Inches(1.0)
    add_rect(slide, cx, cy, Inches(3.4), Inches(0.85), ACCENT_BG, ACCENT_LIGHT)
    add_rect(slide, cx, cy, Pt(5), Inches(0.85), ACCENT)
    add_text(slide, name, cx + Inches(0.2), cy + Inches(0.05), Inches(3.0), Inches(0.35),
             Pt(13), TEXT, bold=True)
    add_text(slide, desc, cx + Inches(0.2), cy + Inches(0.45), Inches(3.0), Inches(0.3),
             Pt(12), MUTED)

# Interaction box
add_rounded_rect(slide, Inches(0.8), Inches(4.35), Inches(7.3), Inches(1.2), BG_LIGHT, BORDER)
add_text(slide, "操作互動機制", Inches(1.1), Inches(4.55), Inches(4), Inches(0.35), Pt(15), TEXT, bold=True)
add_text(slide, "WASD 移動 / 滑鼠視角控制 / 靠近部件按 [E] 觸發「🔍 檢查」「⚙ 操作」選單\nPOST /api/action → 即時送入 SOP 評估引擎",
         Inches(1.1), Inches(4.95), Inches(6.8), Inches(0.55), Pt(13), MUTED, wrap=True)

# Specs panel
spec_groups = [
    ("設備規格（ASML TWINSCAN NXT:870）",
     [("光源波長", "193 nm（ArF）"), ("曝光模式", "浸潤式 DUV"),
      ("模型檔案", "asml_duv.glb"), ("視角", "第一人稱（PointerLock）")]),
    ("兩種訓練模式",
     [("操作模式", "在 3D 執行故障排查 SOP"), ("問答模式", "AI 導師概念性追問"),
      ("觸發點", "關鍵時間點中斷"), ("解鎖方式", "答完方可繼續")]),
]
gy = Inches(2.0)
for gtitle, rows in spec_groups:
    add_rounded_rect(slide, Inches(8.5), gy, Inches(4.5), Inches(2.4), BG_LIGHT, BORDER)
    add_text(slide, gtitle, Inches(8.7), gy + Inches(0.1), Inches(4.1), Inches(0.35),
             Pt(13), ACCENT, bold=True)
    ry = gy + Inches(0.55)
    for k, v in rows:
        add_text(slide, k, Inches(8.7), ry, Inches(1.6), Inches(0.3), Pt(12), MUTED)
        add_text(slide, v, Inches(10.5), ry, Inches(2.3), Inches(0.3), Pt(12), TEXT, bold=True)
        ry += Inches(0.42)
    gy += Inches(2.6)

slide_num(slide, 6)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: SECOM DATA
# ═══════════════════════════════════════════════════════════
slide = add_slide()
top_bar(slide); bottom_bar(slide)

add_label(slide, "核心技術 2", Inches(0.8), Inches(0.55))
add_title(slide, "UCI SECOM 資料集驅動感測器模擬",
          Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.6), Pt(32))

# Stat boxes
for i, (num, lbl) in enumerate([("1,567", "晶圓量測記錄"), ("590", "感測器特徵")]):
    bx = Inches(0.8) + i * Inches(3.3)
    add_rounded_rect(slide, bx, Inches(1.75), Inches(3.0), Inches(1.1), ACCENT, None,
                     num, Pt(32), WHITE)
    add_text(slide, lbl, bx, Inches(2.9), Inches(3.0), Inches(0.3),
             Pt(13), MUTED, align=PP_ALIGN.CENTER)

# Flow steps
flow_steps = [
    ("1", "填補 NaN（各欄中位數）", ""),
    ("2", "PCA 降維（590 → 20 主成分）", ""),
    ("3", "分離 Pass / Fail 樣本建立雜訊模型", "Pass：1,463 筆 ｜ Fail：104 筆"),
    ("4", "計算設備健康分數（health_score 0~1）", "依 PCA 投影座標計算，調變 CD 預測雜訊"),
]
fy = Inches(4.2)
for n, text, sub in flow_steps:
    add_rounded_rect(slide, Inches(0.8), fy, Pt(30), Pt(30), ACCENT, None,
                     n, Pt(12), WHITE)
    add_text(slide, text, Inches(1.5), fy, Inches(5.0), Inches(0.35), Pt(14), TEXT)
    if sub:
        add_text(slide, sub, Inches(1.5), fy + Inches(0.35), Inches(5.0), Inches(0.3), Pt(12), MUTED)
        fy += Inches(0.75)
    else:
        fy += Inches(0.5)

# Right cards
usage_cards = [
    ("SECOM 在本系統中的角色",
     "不用於預測良率，而是建立製程雜訊統計模型，驅動虛擬 HMI 的 590 個感測器即時讀值，呈現符合真實統計分布的設備狀態"),
    ("CD 預測雜訊標準差",
     "正常狀態 σ = 1.5 nm   ／   故障狀態 σ = 3.0–4.8 nm"),
    ("HMI 面板即時呈現",
     "CDU Map 空間分布圖 ／ CD 趨勢圖（含 3σ 控制線）／ Overlay 向量場可視化"),
]
cy7 = Inches(1.7)
for title, body in usage_cards:
    add_rounded_rect(slide, Inches(7.0), cy7, Inches(5.9), Inches(1.5), BG_LIGHT, BORDER)
    add_text(slide, title, Inches(7.25), cy7 + Inches(0.1), Inches(5.4), Inches(0.35),
             Pt(14), ACCENT, bold=True)
    add_text(slide, body, Inches(7.25), cy7 + Inches(0.5), Inches(5.4), Inches(0.9),
             Pt(13), TEXT, wrap=True)
    cy7 += Inches(1.75)

slide_num(slide, 7)

# ═══════════════════════════════════════════════════════════
# SLIDE 8: 3D DEMO MOCKUP
# ═══════════════════════════════════════════════════════════
slide = add_slide()
add_rect(slide, 0, 0, W, H, RGBColor(0x0D, 0x11, 0x17))

add_text(slide, "THREE.JS WEBGL 第一人稱視角",
         Inches(0.5), Inches(0.3), Inches(6), Inches(0.4),
         Pt(13), RGBColor(0x80, 0x80, 0x80), bold=True)

# Machine grid mockup
parts = [
    ("雷射光源\nArF 193nm", False, False),
    ("照明系統\nσ 相干度",  False, False),
    ("真空系統\n真空維持中", False, False),
    ("光罩載台\n對準完成",  False, False),
    ("投影鏡組\n▶ 互動中 [E]", True, False),
    ("⚠ 鏡片熱點\nT +4.2°C", False, True),
    ("晶圓台\n步進中",      False, False),
    ("液浸冷卻\n水溫正常",  False, False),
    ("控制系統\nHMI 連線",  False, False),
]
pc = 0
for row in range(3):
    for col in range(3):
        name, is_active, is_hot = parts[pc]; pc += 1
        px = Inches(1.0) + col * Inches(2.5)
        py = Inches(1.0) + row * Inches(1.8)
        if is_hot:
            bg = RGBColor(0x28, 0x1A, 0x08); lc = RGBColor(0xCC, 0x88, 0x22)
            tc = RGBColor(0xE0, 0xA0, 0x50)
        elif is_active:
            bg = RGBColor(0x18, 0x28, 0x44); lc = RGBColor(0x50, 0x90, 0xCC)
            tc = WHITE
        else:
            bg = RGBColor(0x18, 0x22, 0x30); lc = RGBColor(0x40, 0x60, 0x80)
            tc = RGBColor(0x60, 0x88, 0xAA)
        sh = slide.shapes.add_shape(5, px, py, Inches(2.2), Inches(1.5))
        sh.adjustments[0] = 0.05
        sh.fill.solid(); sh.fill.fore_color.rgb = bg
        sh.line.color.rgb = lc; sh.line.width = Pt(1)
        tf = sh.text_frame; tf.word_wrap = True
        par = tf.paragraphs[0]; par.alignment = PP_ALIGN.CENTER
        r = par.add_run(); r.text = name; r.font.size = Pt(12); r.font.color.rgb = tc

# Caption bar
add_rect(slide, 0, H - Inches(0.55), W - Inches(3.0), Inches(0.55),
         RGBColor(0x00, 0x00, 0x00))
add_text(slide, "故障視覺化：橘色發光 = 鏡片熱點警報  ｜  按 [E] 與零件互動  ｜  按 [C] 開啟 AI 對話",
         Inches(0.3), H - Inches(0.5), W - Inches(3.5), Inches(0.45), Pt(13),
         RGBColor(0xC0, 0xC8, 0xD0))

# HUD right panel
hud_bg = add_rect(slide, W - Inches(3.0), 0, Inches(3.0), H, RGBColor(0x11, 0x18, 0x27))
hud_bg.line.fill.background()
add_text(slide, "HMI 感測器", W - Inches(2.9), Inches(0.3), Inches(2.7), Inches(0.35),
         Pt(12), ACCENT_LIGHT, bold=True)
hud_rows = [("Dose", "32.4 mJ/cm²", None),
            ("Focus", "+0.08 μm ↑", RGBColor(0xE0, 0xA0, 0x50)),
            ("LensTemp", "38.4°C ⚠", RGBColor(0xE0, 0x50, 0x50)),
            ("CDU 3σ", "2.1 nm ⚠",  RGBColor(0xE0, 0x50, 0x50)),
            ("Overlay", "2.8 nm",    None)]
hy = Inches(0.75)
for k, v, vc in hud_rows:
    add_text(slide, k, W - Inches(2.9), hy, Inches(1.3), Inches(0.35), Pt(12), RGBColor(0x80, 0x80, 0x90))
    add_text(slide, v, W - Inches(1.5), hy, Inches(1.3), Inches(0.35), Pt(12), vc or WHITE, bold=True)
    hy += Inches(0.45)

add_text(slide, "自適應模式", W - Inches(2.9), hy + Inches(0.1), Inches(2.7), Inches(0.35),
         Pt(12), ACCENT_LIGHT, bold=True)
add_rounded_rect(slide, W - Inches(2.9), hy + Inches(0.5), Inches(2.6), Inches(0.5),
                 RGBColor(0x18, 0x38, 0x20), RGBColor(0x40, 0x90, 0x55),
                 "🔧 鷹架模式 (Scaffolding)", Pt(12), RGBColor(0x60, 0xC0, 0x80))

slide_num(slide, 8, RGBColor(0x60, 0x60, 0x70))

# ═══════════════════════════════════════════════════════════
# SLIDE 9: AI MENTOR
# ═══════════════════════════════════════════════════════════
slide = add_slide()
top_bar(slide); bottom_bar(slide)

add_label(slide, "核心技術 3 · 圖 3 自適應評分與引導機制圖", Inches(0.8), Inches(0.55), Inches(9))
add_title(slide, "ProactiveMentor AI 導師·雙層語義評估",
          Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.6), Pt(32))

# Coordinator
add_rounded_rect(slide, Inches(0.8), Inches(1.8), Inches(6.2), Inches(0.65), ACCENT, None,
                 "ProactiveMentor 模組", Pt(16), WHITE)

agents = [
    ("主動告警生成", "5 種故障根因×3 個隨機化開場問題變體，避免訓練重覆記憶答案"),
    ("蘇格拉底式追問", "依 ActionSession 即時得分動態調整追問難度"),
    ("同義詞關鍵字群組", "容許「偏移」「shift」「歪掉」等多樣表述，保底 6/7 分或 3 分"),
    ("LLM 語義評估", "回答 / 標準解釋 / 難度組成 prompt，標記 +5/+2/−3 分"),
]
ax = Inches(0.8)
for i, (name, role) in enumerate(agents):
    cx = ax if i % 2 == 0 else ax + Inches(3.2)
    cy = Inches(2.65) if i < 2 else Inches(3.6)
    add_rounded_rect(slide, cx, cy, Inches(3.0), Inches(0.85), ACCENT_BG, ACCENT_LIGHT)
    add_text(slide, name, cx + Inches(0.15), cy + Inches(0.05), Inches(2.7), Inches(0.3),
             Pt(13), ACCENT, bold=True)
    add_text(slide, role, cx + Inches(0.15), cy + Inches(0.4), Inches(2.7), Inches(0.35),
             Pt(11), MUTED, wrap=True)

# LLM card
add_rounded_rect(slide, Inches(0.8), Inches(4.6), Inches(6.2), Inches(0.8), BG_LIGHT, BORDER)
add_text(slide, "本地端 Qwen LLM（通義千問）", Inches(1.0), Inches(4.7), Inches(4), Inches(0.3),
         Pt(14), TEXT, bold=True)
add_text(slide, "晶圓廠資安要求：完全不連網外部 API，本地推理保障製程資料隱私",
         Inches(1.0), Inches(5.05), Inches(5.8), Inches(0.3), Pt(12), MUTED)

# Dialogue demo (right)
add_rounded_rect(slide, Inches(7.3), Inches(1.75), Inches(5.7), Inches(3.6), BG_LIGHT, BORDER)
add_text(slide, "追問示例（故障觸發後主動生成）",
         Inches(7.5), Inches(1.85), Inches(5.2), Inches(0.3), Pt(12), ACCENT, bold=True)
bubbles = [
    ("偵測到鏡片溫度異常，請問你判斷此現象背後的物理機制是什麼？", ACCENT, WHITE),
    ("鏡片吸收雷射能量後熱膨脹，折射率改變導致 Overlay 偏移", RGBColor(0xE5, 0xE7, 0xEB), TEXT),
    ("正確！那麼排查時第一步應該拿什麼讀值來判斷嚴重程度？", ACCENT, WHITE),
    ("查看投影鏡組的溫升量與 CDU 偏差", RGBColor(0xE5, 0xE7, 0xEB), TEXT),
]
by = Inches(2.2)
for btext, bg_c, fg_c in bubbles:
    add_rounded_rect(slide, Inches(7.5), by, Inches(5.2), Inches(0.6), bg_c, None,
                     btext, Pt(12), fg_c)
    by += Inches(0.72)

# Score display
add_rounded_rect(slide, Inches(7.3), Inches(5.5), Inches(5.7), Inches(1.1), ACCENT, None)
add_text(slide, "+5", Inches(7.5), Inches(5.6), Inches(0.9), Inches(0.9), Pt(42), WHITE, bold=True)
add_text(slide, "概念理解、所有同義詞群組命中\n→ LLM 語義確認·加 5 分\n→ ActionSession 即時在 100 分制上提升",
         Inches(8.5), Inches(5.6), Inches(4.3), Inches(0.9), Pt(12), WHITE)

slide_num(slide, 9)

# ═══════════════════════════════════════════════════════════
# SLIDE 10: ADAPTIVE TEACHING
# ═══════════════════════════════════════════════════════════
slide = add_slide()
top_bar(slide); bottom_bar(slide)

add_label(slide, "核心技術 4 · 表 2 提示層級與分數對應", Inches(0.8), Inches(0.55), Inches(9))
add_title(slide, "自適應四模式 · 依 ActionSession 得分動態切換",
          Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.6), Pt(32))

modes = [
    ("挑戰\nChallenge",   "得分 > 85 分", "不給具體提示，引導自行分析警報；追問深入原理、量化估算",
     RGBColor(0x18, 0x28, 0x3E), WHITE, "> 85"),
    ("標準\nStandard",    "得分 64–85 分", "先說明讀值偏差，再指出應前往的部件；確認重點概念",
     ACCENT_BG, TEXT, "64–85"),
    ("鷹架\nScaffolding", "得分 40–63 分", "給出部件名稱 + 操作動作方向；提供方向性引導",
     RGBColor(0xFE, 0xF9, 0xEE), TEXT, "40–63"),
    ("補救\nRemedial",    "得分 < 40 分",  "完整步驟：部件、動作、原因全說明；解釋基本定義",
     RGBColor(0xFE, 0xF0, 0xEE), TEXT, "< 40"),
]
mx = Inches(0.8)
for name, trigger, style, bg_c, tc, score in modes:
    sh = slide.shapes.add_shape(5, mx, Inches(1.85), Inches(2.9), Inches(5.0))
    sh.adjustments[0] = 0.04
    sh.fill.solid(); sh.fill.fore_color.rgb = bg_c
    sh.line.fill.background()
    add_text(slide, name, mx + Inches(0.2), Inches(2.0), Inches(2.5), Inches(0.75),
             Pt(22), tc, bold=True)
    add_rounded_rect(slide, mx + Inches(0.15), Inches(2.9), Inches(2.6), Inches(0.6),
                     RGBColor(0, 0, 0), None, trigger, Pt(12),
                     tc if bg_c == RGBColor(0x18, 0x28, 0x3E) else MUTED)
    add_text(slide, style, mx + Inches(0.2), Inches(3.65), Inches(2.55), Inches(1.5),
             Pt(13), tc, wrap=True)
    add_text(slide, score, mx + Inches(0.2), Inches(6.0), Inches(2.5), Inches(0.5),
             Pt(26), tc, bold=True)
    mx += Inches(3.1)

add_text(slide, "起始 100 分 · 錯誤操作依部件/動作匹配程度扣 2–15 分 · 求助學長扣 10 分後給差異化提示 · 操作評分與概念追問難度由同一得分驅動",
         Inches(0.8), Inches(6.9), Inches(11.8), Inches(0.45), Pt(13), MUTED, wrap=True)

slide_num(slide, 10)

# ═══════════════════════════════════════════════════════════
# SLIDE 11: FAULT SCENARIOS
# ═══════════════════════════════════════════════════════════
slide = add_slide()
top_bar(slide); bottom_bar(slide)

add_label(slide, "訓練設計 · 表 1 故障情境與 SOP 步驟數", Inches(0.8), Inches(0.55), Inches(9))
add_title(slide, "五種故障情境 × SOP 自主評估",
          Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.6), Pt(32))

faults = [
    ("01", "投影鏡片過熱（lens_hotspot） · 5 步", "鏡片過熱→折射率變化→Overlay 誤差增加", True),
    ("02", "光罩污染（contamination） · 5 步", "光源強度下降，卸載光罩→目視檢查→清潔/更換→重新對準", False),
    ("03", "晶圓載台位置誤差（stage_error） · 4 步", "台位移→Overlay 超差→校正定位精度→驗證", False),
    ("04", "曝光 Dose 漂移（dose_drift） · 4 步", "雷射能量監測器異常→CDU 惡化", False),
    ("05", "焦距異常漂移（focus_drift） · 4 步", "AF 感測器異常→Focus Map 偏移→CD 失控", False),
]
fy = Inches(1.85)
for idx, name, desc, active in faults:
    bg_c = ACCENT_BG if active else BG_LIGHT
    lc = ACCENT if active else BORDER
    add_rounded_rect(slide, Inches(0.8), fy, Inches(6.0), Inches(0.88), bg_c, lc)
    add_text(slide, idx, Inches(1.0), fy + Inches(0.08), Inches(0.7), Inches(0.6),
             Pt(22), ACCENT_LIGHT, bold=True)
    add_text(slide, name, Inches(1.8), fy + Inches(0.08), Inches(4.8), Inches(0.35),
             Pt(13), TEXT, bold=True)
    add_text(slide, desc, Inches(1.8), fy + Inches(0.47), Inches(4.8), Inches(0.35),
             Pt(12), MUTED)
    fy += Inches(1.02)

# SOP steps
add_text(slide, "情境 01：鏡片熱點 SOP（5 步驟）",
         Inches(7.2), Inches(1.75), Inches(5.8), Inches(0.4), Pt(14), ACCENT, bold=True)
sop_steps = [
    ("1", "檢查投影鏡組溫升量", "確認嚴重程度，判斷後續處置"),
    ("2", "降低 Dose 設定", "減少熱輸入，阻止溫升惡化"),
    ("3", "停止曝光，等待冷卻", "讓鏡片自然冷卻至規格內"),
    ("4", "確認液浸冷卻水迴路", "查看冷卻水流量與水溫是否正常"),
    ("5", "恢復曝光，監控 CDU", "驗證 CDU 回到規格內"),
]
sy = Inches(2.25)
for n, action, reason in sop_steps:
    add_rounded_rect(slide, Inches(7.2), sy, Inches(5.8), Inches(0.82), ACCENT_BG, None)
    add_rounded_rect(slide, Inches(7.2), sy + Inches(0.2), Pt(32), Pt(32), ACCENT, None,
                     n, Pt(12), WHITE)
    add_text(slide, action, Inches(7.85), sy + Inches(0.05), Inches(4.9), Inches(0.3),
             Pt(13), TEXT, bold=True)
    add_text(slide, reason, Inches(7.85), sy + Inches(0.42), Inches(4.9), Inches(0.3),
             Pt(12), MUTED)
    sy += Inches(0.93)

add_rounded_rect(slide, Inches(7.2), Inches(6.95), Inches(5.8), Inches(0.6),
                 RGBColor(0xEE, 0xF4, 0xFA), ACCENT_LIGHT,
                 "⚠ 自主評估：學員自行判斷操作順序。錯誤分三級，結合自適應模式扣 2–15 分；求助學長扣 10 分",
                 Pt(11), ACCENT)

slide_num(slide, 11)

# ═══════════════════════════════════════════════════════════
# SLIDE 12: SCORING
# ═══════════════════════════════════════════════════════════
slide = add_slide()
top_bar(slide); bottom_bar(slide)

add_label(slide, "評量設計 · 表 2 提示層級與分數對應", Inches(0.8), Inches(0.55), Inches(9))
add_title(slide, "100 分制 ActionSession 評分 · 四級自適應提示",
          Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.6), Pt(32))

score_cards = [
    ("100", "起始分數", "ActionSession 初始 100 分，依操作動態加減"),
    ("−2∼15", "錯誤扣分", "partial_action / partial_component / full_wrong × 自適應模式"),
    ("−10", "求助學長", "詢問學長扣 10 分後，給予當前分數對應的差異化提示"),
]
sc_y = Inches(1.85)
for pct, name, desc in score_cards:
    add_rounded_rect(slide, Inches(0.8), sc_y, Inches(6.0), Inches(1.2), BG_LIGHT, BORDER)
    add_text(slide, pct, Inches(1.0), sc_y + Inches(0.1), Inches(1.5), Inches(0.85),
             Pt(32), ACCENT, bold=True)
    add_text(slide, name, Inches(2.7), sc_y + Inches(0.1), Inches(3.9), Inches(0.35),
             Pt(15), TEXT, bold=True)
    add_text(slide, desc, Inches(2.7), sc_y + Inches(0.55), Inches(3.9), Inches(0.55),
             Pt(12), MUTED, wrap=True)
    sc_y += Inches(1.35)

# Bonus card
add_rounded_rect(slide, Inches(0.8), sc_y + Inches(0.15), Inches(6.0), Inches(0.9),
                 RGBColor(0xFE, 0xF9, 0xEE), RGBColor(0xCC, 0xA0, 0x40))
add_text(slide, "⊕ 雙層語義評估", Inches(1.0), sc_y + Inches(0.25), Inches(2.5), Inches(0.35),
         Pt(14), RGBColor(0x80, 0x60, 0x10), bold=True)
add_text(slide, "同義詞群組 +7/+3  ·  LLM 語義評估 +5/+2/−3",
         Inches(3.7), sc_y + Inches(0.3), Inches(3.0), Inches(0.35), Pt(13), RGBColor(0x70, 0x50, 0x10))

# Grade right
grades = [
    ("★", "> 85 分",  "挑戰 Challenge  不給具體提示，引導自行分析警報", ACCENT_BG, ACCENT),
    ("◎", "64–85 分", "標準 Standard  先看讀值偏差，再指出應前往部件", ACCENT_BG, ACCENT),
    ("▲", "40–63 分", "鷹架 Scaffolding  給部件名稱 + 操作動作方向", BG_LIGHT, RGBColor(0x80, 0x60, 0x10)),
    ("↩", "< 40 分",  "補救 Remedial  完整步驟：部件、動作、原因全說明", BG_LIGHT, RGBColor(0x80, 0x25, 0x15)),
]
gy = Inches(1.85)
for letter, rng, desc, bg_c, lc in grades:
    add_rounded_rect(slide, Inches(7.3), gy, Inches(5.7), Inches(1.1), bg_c, lc)
    add_text(slide, letter, Inches(7.5), gy + Inches(0.1), Inches(0.7), Inches(0.85),
             Pt(34), lc, bold=True)
    add_text(slide, rng, Inches(8.3), gy + Inches(0.1), Inches(4.4), Inches(0.35),
             Pt(15), TEXT, bold=True)
    add_text(slide, desc, Inches(8.3), gy + Inches(0.5), Inches(4.4), Inches(0.5),
             Pt(12), MUTED, wrap=True)
    gy += Inches(1.3)

slide_num(slide, 12)

# ═══════════════════════════════════════════════════════════
# SLIDE 13: EXPERIMENT RESULTS
# ═══════════════════════════════════════════════════════════
slide = add_slide()
top_bar(slide); bottom_bar(slide)

add_label(slide, "實驗結果 · 表 4 前測 / 後測進步統計表", Inches(0.8), Inches(0.55), Inches(9))
add_title(slide, "訓練成效評估 · 前測 10.45 → 後測 11.55 分",
          Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.6), Pt(32))

# Left stat box
add_rounded_rect(slide, Inches(0.8), Inches(1.85), Inches(4.5), Inches(1.6), ACCENT, None)
add_text(slide, "整體平均進步幅度", Inches(1.0), Inches(1.95), Inches(4.1), Inches(0.35),
         Pt(13), WHITE)
add_text(slide, "+7.27%", Inches(1.0), Inches(2.35), Inches(4.1), Inches(0.85),
         Pt(42), WHITE, bold=True)
add_text(slide, "滿分 15 分，平均進步 +1.09 分",
         Inches(1.0), Inches(3.2), Inches(4.1), Inches(0.3), Pt(13), WHITE)

add_rounded_rect(slide, Inches(0.8), Inches(3.6), Inches(4.5), Inches(1.1), BG_LIGHT, BORDER)
add_text(slide, "實驗設計", Inches(1.0), Inches(3.7), Inches(4.1), Inches(0.3), Pt(12), MUTED)
add_text(slide, "前測 → 系統操作（20–30 分）→ 後測\n概念理解問卷 15 分滿分",
         Inches(1.0), Inches(4.05), Inches(4.1), Inches(0.55), Pt(13), TEXT)

add_rounded_rect(slide, Inches(0.8), Inches(4.85), Inches(4.5), Inches(1.1), BG_LIGHT, BORDER)
add_text(slide, "受試者分組", Inches(1.0), Inches(4.95), Inches(4.1), Inches(0.3), Pt(12), MUTED)
add_text(slide, "非半導體業 ／ 工程師（0–5 年）／ 資深工程師（5 年以上）",
         Inches(1.0), Inches(5.3), Inches(4.1), Inches(0.55), Pt(13), TEXT)

# Data table
table_left = Inches(5.7)
table_top  = Inches(1.85)
headers = ["職業背景", "前測", "後測", "進步分數", "進步幅度"]
col_ws = [Inches(2.4), Inches(0.9), Inches(0.9), Inches(1.0), Inches(1.0)]

# Header
hx = table_left
add_rect(slide, table_left, table_top, sum(col_ws), Inches(0.55), ACCENT)
for i, h in enumerate(headers):
    add_text(slide, h, hx + Inches(0.1), table_top + Inches(0.1), col_ws[i] - Inches(0.1),
             Inches(0.35), Pt(13), WHITE, bold=True)
    hx += col_ws[i]

rows_data = [
    ("非半導體業人員",          "8.83", "9.67",  "+0.83", "5.53%",  False),
    ("半導體工程師（0–5 年）",   "11.67","13.33", "+1.67", "11.13%", True),
    ("資深工程師（5 年以上）",   "13.50","14.50", "+1.00", "6.67%",  False),
    ("整體平均成效",            "10.45","11.55", "+1.09", "7.27%",  False, True),
]
ry = table_top + Inches(0.55)
for i, rd in enumerate(rows_data):
    is_highlight = len(rd) > 6 and rd[6]
    bg_c = ACCENT if is_highlight else (ACCENT_BG if rd[5] else (BG_LIGHT if i % 2 == 0 else WHITE))
    tc = WHITE if is_highlight else TEXT
    add_rect(slide, table_left, ry, sum(col_ws), Inches(0.7), bg_c, BORDER)
    rx = table_left
    for j, cell in enumerate(rd[:5]):
        is_good = j >= 3
        cell_c = WHITE if is_highlight else (RGBColor(0x20, 0x80, 0x48) if is_good and not is_highlight else tc)
        add_text(slide, cell, rx + Inches(0.1), ry + Inches(0.18), col_ws[j] - Inches(0.1),
                 Inches(0.35), Pt(13), cell_c, bold=(j == 0 or is_highlight))
        rx += col_ws[j]
    ry += Inches(0.7)

add_text(slide, "最高成效：0–5 年新進工程師進步 11.13%，證明系統對具基礎的新進人員有實戰銜接效果",
         table_left, ry + Inches(0.1), sum(col_ws), Inches(0.45), Pt(13), TEXT, wrap=True)

slide_num(slide, 13)

# ═══════════════════════════════════════════════════════════
# SLIDE 14: CONTRIBUTIONS
# ═══════════════════════════════════════════════════════════
slide = add_slide()
top_bar(slide); bottom_bar(slide)

add_label(slide, "學術貢獻", Inches(0.8), Inches(0.55))
add_title(slide, "四項主要貢獻與未來展望",
          Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.6), Pt(32))

contribs = [
    ("01", "WebGL 數位孿生環境",
     "以 WebGL 建置 ArF 曝光機數位孿生 3D 訓練環境，支援第一人稱沉浸式操作互動，瀏覽器即可部署",
     "→ 可延伸至 ALD / PVD / CMP"),
    ("02", "自適應 SOP 評估引擎",
     "設計 0～100 分制自適應評分引擎，動態調整回饋層級與提示詳細程度，整合操作評分與提示強度",
     "→ 擴充更多故障情境類型"),
    ("03", "蘇格拉底式 AI 追問機制",
     "結合同義詞關鍵字群組與 LLM，對概念型與行動型問題分別進行差異化評估，解決表述多樣性問題",
     "→ 多人協作訓練模式"),
    ("04", "ProactiveMentor 整合",
     "故障發生時主動生成告警與情境化問題，依學員得分自動調整追問難度；本地端 Qwen 符合晶圓廠資安要求",
     "→ 擴大受試者樣本進一步驗證"),
]
cx_list = [Inches(0.8), Inches(6.9)]
cy_list = [Inches(1.9), Inches(4.6)]
for i, (num, title, body, future) in enumerate(contribs):
    cx = cx_list[i % 2]
    cy = cy_list[i // 2]
    add_rounded_rect(slide, cx, cy, Inches(5.8), Inches(2.5), BG_LIGHT, BORDER)
    add_text(slide, num, cx + Inches(0.25), cy + Inches(0.15), Inches(1.0), Inches(0.6),
             Pt(28), ACCENT_LIGHT, bold=True)
    add_text(slide, title, cx + Inches(0.25), cy + Inches(0.75), Inches(5.2), Inches(0.4),
             Pt(16), TEXT, bold=True)
    add_text(slide, body, cx + Inches(0.25), cy + Inches(1.2), Inches(5.2), Inches(0.85),
             Pt(13), MUTED, wrap=True)
    add_rounded_rect(slide, cx + Inches(0.2), cy + Inches(2.1), Inches(4.5), Inches(0.35),
                     ACCENT_BG, None, future, Pt(12), ACCENT)

slide_num(slide, 14)

# ═══════════════════════════════════════════════════════════
# SLIDE 15: CLOSING / Q&A
# ═══════════════════════════════════════════════════════════
slide = add_slide()
add_rect(slide, 0, 0, W, H, ACCENT)

add_title(slide, "感謝聆聽",
          Inches(2.0), Inches(1.5), Inches(9.5), Inches(1.4), Pt(60), WHITE)
add_text(slide, "半導體微影設備虛擬訓練系統\n基於數位孿生與自適應 AI 導師",
         Inches(2.0), Inches(3.1), Inches(9.5), Inches(1.0), Pt(22), WHITE, align=PP_ALIGN.CENTER)

# Meta
meta_items = [("研究生", "趙威豪 M1321114"), ("學校", "長庚大學 碩士班"), ("日期", "2026 年 4 月")]
mx_start = Inches(2.0)
for label, val in meta_items:
    add_text(slide, label, mx_start, Inches(4.3), Inches(2.5), Inches(0.35), Pt(13),
             RGBColor(0xC0, 0xD0, 0xE0), align=PP_ALIGN.CENTER)
    add_text(slide, val, mx_start, Inches(4.7), Inches(2.5), Inches(0.4), Pt(16), WHITE,
             bold=True, align=PP_ALIGN.CENTER)
    mx_start += Inches(3.2)

# Q&A box
add_rounded_rect(slide, Inches(4.5), Inches(5.5), Inches(4.5), Inches(1.0),
                 RGBColor(0x18, 0x50, 0x7A), RGBColor(0x60, 0xA0, 0xCC),
                 "Q & A", Pt(26), WHITE)

slide_num(slide, 15, RGBColor(0xFF, 0xFF, 0xFF))

# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
out_path = r"c:\Users\user\Desktop\在職碩\OneDrive - 長庚大學\長庚碩班\論文\口試報告.pptx"
prs.save(out_path)
print("Saved: " + out_path)
