"""
自動生成自適應教學系統簡報
根據 自適應判斷-PPT版流程圖.md 生成 PowerPoint 投影片
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_adaptive_teaching_ppt():
    """創建自適應教學系統簡報"""

    # 創建簡報
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 定義顏色方案
    COLORS = {
        'challenge': RGBColor(255, 87, 51),    # 紅橙色（挑戰模式）
        'standard': RGBColor(52, 152, 219),    # 藍色（標準模式）
        'scaffolding': RGBColor(241, 196, 15), # 黃色（鷹架模式）
        'remedial': RGBColor(46, 204, 113),    # 綠色（補救模式）
        'title': RGBColor(44, 62, 80),         # 深灰色
        'text': RGBColor(52, 73, 94)           # 灰色
    }

    # ========== Slide 1: 核心判斷邏輯 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面

    # 標題
    title = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title.text_frame
    title_frame.text = "核心判斷邏輯（一頁總覽）"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['title']
    title_para.alignment = PP_ALIGN.CENTER

    # 流程圖文字
    flowchart = slide1.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(5))
    tf = flowchart.text_frame
    tf.word_wrap = True

    flow_text = """    學員回答
        ↓
    AI 評分 (0-10分)
        ↓
    記錄歷史 + 計算平均
        ↓
    平均 ≥ 8.5分？
        ↓
    YES → 🎯 挑戰模式
    NO  → 平均 ≥ 6.5分？
          YES → ✅ 標準模式
          NO  → 平均 ≥ 4.0分？
                YES → 🏗️ 鷹架模式
                NO  → 🔧 補救模式"""

    tf.text = flow_text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.name = 'Consolas'
        paragraph.font.color.rgb = COLORS['text']
        paragraph.alignment = PP_ALIGN.CENTER

    # ========== Slide 2: 4種模式對比表 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title2_frame = title2.text_frame
    title2_frame.text = "4種教學模式對比"
    title2_para = title2_frame.paragraphs[0]
    title2_para.font.size = Pt(32)
    title2_para.font.bold = True
    title2_para.font.color.rgb = COLORS['title']
    title2_para.alignment = PP_ALIGN.CENTER

    # 創建表格
    rows, cols = 6, 5
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(8.4)
    height = Inches(5)

    table = slide2.shapes.add_table(rows, cols, left, top, width, height).table

    # 設定表格內容
    table_data = [
        ['模式', '🎯 CHALLENGE', '✅ STANDARD', '🏗️ SCAFFOLDING', '🔧 REMEDIAL'],
        ['觸發條件', '≥8.5分', '6.5-8.5', '4.0-6.5', '<4.0'],
        ['問題難度', '複雜整合', '適中', '基礎', '最簡單'],
        ['提示', '無', '少量', '頻繁', '大量'],
        ['反饋', '簡潔', '平衡', '詳細', '超詳細'],
        ['適合學員', '高手', '一般', '初學者', '掙扎中']
    ]

    # 填充表格
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text

            # 設定文字格式
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.bold = (row_idx == 0)  # 標題列加粗

                # 標題列顏色
                if row_idx == 0:
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)

            # 設定背景顏色
            if row_idx == 0:
                if col_idx == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS['challenge']
                elif col_idx == 2:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS['standard']
                elif col_idx == 3:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS['scaffolding']
                elif col_idx == 4:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS['remedial']
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS['title']

    # ========== Slide 3: 評估維度 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title3_frame = title3.text_frame
    title3_frame.text = "評估維度（多維度評分）"
    title3_para = title3_frame.paragraphs[0]
    title3_para.font.size = Pt(32)
    title3_para.font.bold = True
    title3_para.font.color.rgb = COLORS['title']
    title3_para.alignment = PP_ALIGN.CENTER

    # 評估維度圖
    dimensions = slide3.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(3))
    dim_tf = dimensions.text_frame
    dim_tf.text = """        最終決策
           ▲
           │
    ┌──────┼──────┐
    │      │      │
   40%    30%    20%    10%
    │      │      │      │
最近5次  單次   主題   學習
 平均    得分   掌握   趨勢"""

    for paragraph in dim_tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.name = 'Consolas'
        paragraph.alignment = PP_ALIGN.CENTER

    # 說明文字
    note = slide3.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1))
    note_tf = note.text_frame
    note_tf.text = "不只看一次表現，綜合評估！"
    note_para = note_tf.paragraphs[0]
    note_para.font.size = Pt(24)
    note_para.font.bold = True
    note_para.font.color.rgb = COLORS['challenge']
    note_para.alignment = PP_ALIGN.CENTER

    # ========== Slide 4: 學習軌跡範例 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title4_frame = title4.text_frame
    title4_frame.text = "學習軌跡範例"
    title4_para = title4_frame.paragraphs[0]
    title4_para.font.size = Pt(32)
    title4_para.font.bold = True
    title4_para.font.color.rgb = COLORS['title']
    title4_para.alignment = PP_ALIGN.CENTER

    # 軌跡圖
    trajectory = slide4.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(4))
    traj_tf = trajectory.text_frame
    traj_tf.text = """得分
10│              ●━●  [挑戰模式]
 9│         ●━━━╯
 8│    ●━━━╯      [標準模式]
 7│   ╱
 6│  ●            [鷹架模式]
 5│   ●
  └───────────→ 時間
    Q1 Q2 Q3 Q4 Q5 Q6"""

    for paragraph in traj_tf.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.name = 'Consolas'
        paragraph.alignment = PP_ALIGN.CENTER

    # 說明
    note4 = slide4.shapes.add_textbox(Inches(2), Inches(6), Inches(6), Inches(0.8))
    note4_tf = note4.text_frame
    note4_tf.text = "從掙扎到掌握，系統自動調整難度"
    note4_para = note4_tf.paragraphs[0]
    note4_para.font.size = Pt(22)
    note4_para.font.bold = True
    note4_para.font.color.rgb = COLORS['standard']
    note4_para.alignment = PP_ALIGN.CENTER

    # ========== Slide 5: 主題追蹤 ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title5_frame = title5.text_frame
    title5_frame.text = "主題追蹤（知識盲點識別）"
    title5_para = title5_frame.paragraphs[0]
    title5_para.font.size = Pt(32)
    title5_para.font.bold = True
    title5_para.font.color.rgb = COLORS['title']
    title5_para.alignment = PP_ALIGN.CENTER

    # 主題表現
    topics = slide5.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
    topics_tf = topics.text_frame
    topics_tf.text = """學員A的主題表現：

✅ thermal (熱學):  8.5分 - 已掌握
✅ optical (光學):  7.8分 - 精通中
🟡 mechanical:      6.0分 - 發展中
⚠️ vacuum (真空):   4.2分 - 需加強
⚠️ chemical (化學): 3.5分 - 需加強

→ 系統自動針對弱點主題加強訓練"""

    for paragraph in topics_tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.alignment = PP_ALIGN.LEFT

    # ========== Slide 6: 模式切換範例 ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title6_frame = title6.text_frame
    title6_frame.text = "模式切換範例：學習「熱膨脹」"
    title6_para = title6_frame.paragraphs[0]
    title6_para.font.size = Pt(30)
    title6_para.font.bold = True
    title6_para.font.color.rgb = COLORS['title']
    title6_para.alignment = PP_ALIGN.CENTER

    # 切換過程
    switching = slide6.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    switch_tf = switching.text_frame
    switch_tf.text = """第1-2次：得分 6, 5 分
  → 平均 5.5 分
  → 切換至 🏗️ 鷹架模式
  → 系統：分解步驟，逐步引導

第3-5次：得分 7, 8, 9 分
  → 平均 7.0 分
  → 切換至 ✅ 標準模式
  → 系統：正常難度問題

第6-7次：得分 9, 9 分
  → 平均 8.6 分 (≥8.5)
  → 切換至 🎯 挑戰模式
  → 系統：複雜整合問題，無提示"""

    for paragraph in switch_tf.paragraphs:
        paragraph.font.size = Pt(18)

    # ========== Slide 7: 類比說明 ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title7_frame = title7.text_frame
    title7_frame.text = "類比：心理診斷 AI"
    title7_para = title7_frame.paragraphs[0]
    title7_para.font.size = Pt(32)
    title7_para.font.bold = True
    title7_para.font.color.rgb = COLORS['title']
    title7_para.alignment = PP_ALIGN.CENTER

    # 類比內容
    analogy = slide7.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    analogy_tf = analogy.text_frame
    analogy_tf.text = """心理 AI：
  病人憂鬱程度 → 調整問診方式
  輕度 → 衛教建議
  中度 → 心理諮商
  重度 → 精神科轉介

訓練 AI：
  學員理解程度 → 調整教學方式
  優異 → 挑戰模式 (複雜問題)
  良好 → 標準模式 (適中難度)
  普通 → 鷹架模式 (分解步驟)
  不佳 → 補救模式 (基礎重建)"""

    for paragraph in analogy_tf.paragraphs:
        paragraph.font.size = Pt(18)

    # ========== Slide 8: 系統優勢 ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title8_frame = title8.text_frame
    title8_frame.text = "傳統系統 vs 自適應系統"
    title8_para = title8_frame.paragraphs[0]
    title8_para.font.size = Pt(32)
    title8_para.font.bold = True
    title8_para.font.color.rgb = COLORS['title']
    title8_para.alignment = PP_ALIGN.CENTER

    # 對比表格
    rows2, cols2 = 5, 3
    table2 = slide8.shapes.add_table(rows2, cols2, Inches(1.5), Inches(2), Inches(7), Inches(3.5)).table

    comparison_data = [
        ['項目', '傳統', '自適應'],
        ['難度', '固定', '動態調整 ✅'],
        ['個人化', '❌', '✅'],
        ['弱點識別', '人工', '自動 ✅'],
        ['學習動機', '低', '高 ✅']
    ]

    for row_idx, row_data in enumerate(comparison_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table2.cell(row_idx, col_idx)
            cell.text = cell_text

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.bold = (row_idx == 0)

                if row_idx == 0:
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['title']

    # 結果說明
    result = slide8.shapes.add_textbox(Inches(2), Inches(6), Inches(6), Inches(1))
    result_tf = result.text_frame
    result_tf.text = "結果：高手不無聊 ✅ | 弱者不挫折 ✅"
    result_para = result_tf.paragraphs[0]
    result_para.font.size = Pt(20)
    result_para.font.bold = True
    result_para.font.color.rgb = COLORS['challenge']
    result_para.alignment = PP_ALIGN.CENTER

    # ========== Slide 9: 量化成果 ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title9_frame = title9.text_frame
    title9_frame.text = "量化成果"
    title9_para = title9_frame.paragraphs[0]
    title9_para.font.size = Pt(32)
    title9_para.font.bold = True
    title9_para.font.color.rgb = COLORS['title']
    title9_para.alignment = PP_ALIGN.CENTER

    # 量化數據
    metrics = slide9.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(4))
    metrics_tf = metrics.text_frame
    metrics_tf.text = """評估維度      傳統    自適應
━━━━━━━━━━━━━━━━━━━━━━━
難度等級        1       4種
評估維度        1       4維
主題追蹤       無       6類
模式切換       無       自動
個人化        低       高"""

    for paragraph in metrics_tf.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.name = 'Consolas'
        paragraph.alignment = PP_ALIGN.CENTER

    # ========== Slide 10: 完整流程 ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title10_frame = title10.text_frame
    title10_frame.text = "完整流程（簡化版）"
    title10_para = title10_frame.paragraphs[0]
    title10_para.font.size = Pt(32)
    title10_para.font.bold = True
    title10_para.font.color.rgb = COLORS['title']
    title10_para.alignment = PP_ALIGN.CENTER

    # 完整流程
    full_flow = slide10.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(5))
    flow_tf = full_flow.text_frame
    flow_tf.text = """┌─────────┐
│ 學員回答 │
└────┬────┘
     ↓
┌─────────┐
│ AI 評分  │
└────┬────┘
     ↓
┌─────────────┐
│ 多維度評估   │
│ • 最近5次    │
│ • 單次得分   │
│ • 主題分析   │
└────┬────────┘
     ↓
┌─────────────┐
│ 判斷理解程度 │
└────┬────────┘
     ↓
┌─────────────┐
│ 切換教學模式 │
└────┬────────┘
     ↓
┌─────────────┐
│ 生成下一問題 │
└─────────────┘"""

    for paragraph in flow_tf.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.name = 'Consolas'
        paragraph.alignment = PP_ALIGN.CENTER

    # 儲存檔案
    output_path = "自適應教學系統簡報.pptx"
    prs.save(output_path)
    print("[OK] Presentation created: " + output_path)
    print("[OK] Total slides: 10")
    print("[OK] Estimated time: 7-9 minutes")

    return output_path

if __name__ == "__main__":
    print("=" * 60)
    print("Adaptive Teaching System Presentation Generator")
    print("=" * 60)
    print()

    output = create_adaptive_teaching_ppt()

    print()
    print("=" * 60)
    print("Usage:")
    print("1. Open with PowerPoint")
    print("2. Adjust fonts, colors, animations as needed")
    print("3. Refer to presentation script in markdown files")
    print("=" * 60)
